from dotenv import load_dotenv
import os
from openai import OpenAI
from pptx import Presentation
from docx import Document

load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

def main():
    print("\nAI Sales Toolkit Generator\n")

    company = "Hudson Financial Group"
    website = "https://www.hfgventures.net/"

    print("\nGenerating full sales system... please wait...\n")

    # ---------------------------
    # PRESENTATION REQUEST
    # ---------------------------
    presentation_prompt = f"""
You are an elite AI consultant and sales strategist.

Create CLIENT-FACING PRESENTATION CONTENT for this company.

Company: {company}
Website: {website}

STYLE:
- aggressive
- persuasive
- sleek
- modern
- bold
- no fluff
- bullets + short supporting lines
- focused on revenue, competition, missed opportunity, urgency

IMPORTANT:
- Small to mid-size business logic only
- No enterprise-company framing
- Keep it relatable to a reachable decision-maker
- Do not explain implementation steps
- Focus on outcomes and business value

Return EXACTLY 10 slides in this format:

SLIDE 1 TITLE: ...
SLIDE 1 BULLETS:
- ...
- ...
- ...

SLIDE 2 TITLE: ...
SLIDE 2 BULLETS:
- ...
- ...
- ...

Continue this exact pattern through SLIDE 10.

Use these slide themes:

1. Title / AI Opportunity
2. 5 High-Impact AI Facts for the Industry
3. Why This Matters Now
4. Revenue and Efficiency Projections (low / mid / high)
5. AI Opportunities - Sales
6. AI Opportunities - Marketing
7. AI Opportunities - Operations
8. What Competitors Are Doing
9. What We Recommend First
10. Call to Action
"""

    presentation_response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": presentation_prompt}]
    )

    presentation_text = presentation_response.choices[0].message.content

    # ---------------------------
    # TOOLKIT REQUEST
    # ---------------------------
    toolkit_prompt = f"""
You are an elite AI consultant and sales strategist.

Create an INTERNAL SALES TOOLKIT for this company.

Company: {company}
Website: {website}

IMPORTANT:
- Optimize for a solo entrepreneur selling AI services
- Target small to mid-size businesses only
- No enterprise-company nonsense
- Keep recommendations realistic, sellable, and commercially sharp
- Use strong human business judgment
- Focus on what will help close this company and similar companies

Include:

1. What to Sell First
2. Pitch Strategy
3. Objections and Responses
4. Personal Text Message (sounds human, not AI)
5. Cold Email Sequence
   - Email 1
   - Follow-up 1
   - Follow-up 2
6. Similar Companies to Target
   - same or similar industry
   - same metroplex first
   - then expand outward
   - grouped by approximate size
   - only reachable businesses

Write it like a polished sales playbook.
"""

    toolkit_response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": toolkit_prompt}]
    )

    toolkit_text = toolkit_response.choices[0].message.content

    # ---------------------------
    # CREATE POWERPOINT
    # ---------------------------
    prs = Presentation()

    def extract_slide(text, slide_num):
        title_marker = f"SLIDE {slide_num} TITLE:"
        bullets_marker = f"SLIDE {slide_num} BULLETS:"
        next_title_marker = f"SLIDE {slide_num + 1} TITLE:"

        start_title = text.find(title_marker)
        start_bullets = text.find(bullets_marker)

        if start_title == -1 or start_bullets == -1:
            return f"Slide {slide_num}", "Content could not be parsed cleanly."

        title = text[start_title + len(title_marker):start_bullets].strip()

        if slide_num < 10:
            end_bullets = text.find(next_title_marker)
            bullets = text[start_bullets + len(bullets_marker):end_bullets].strip()
        else:
            bullets = text[start_bullets + len(bullets_marker):].strip()

        return title, bullets

    for i in range(1, 11):
        title, bullets = extract_slide(presentation_text, i)

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = bullets

    prs.save("hudson_presentation.pptx")

    # ---------------------------
    # CREATE WORD DOC
    # ---------------------------
    doc = Document()
    doc.add_heading("Hudson Financial Group - AI Sales Toolkit", 0)
    doc.add_paragraph(toolkit_text)
    doc.save("hudson_sales_toolkit.docx")

    print("\nCreated:")
    print("hudson_presentation.pptx")
    print("hudson_sales_toolkit.docx")

if __name__ == "__main__":
    main()